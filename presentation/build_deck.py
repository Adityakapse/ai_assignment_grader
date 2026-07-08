"""
Build the dissertation defence deck (PowerPoint).

Simple, clean template. Structure:
  - no footer
  - a title-only divider slide before each section
  - each research question on its own slide, then its graphs (one graph per slide)
  - per-model graphs (RQ1, RQ4) shown Qwen vs gpt-oss side-by-side
  - real graph PNGs embedded from graph_store/

Run:  python presentation/build_deck.py
"""

import os
import struct

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- config
PROJECT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
G = os.path.join(PROJECT, "graph_store")
QW = os.path.join(G, "Qwen3b", "overall")
GO = os.path.join(G, "gpt-oss-120b", "overall")
MC = os.path.join(G, "Model Comparison", "overall")
ARCH = os.path.join(PROJECT, "architecture_pipeline.png")
OUT = os.path.join(os.path.dirname(__file__), "LLM_Grading_Dissertation.pptx")

FONT = "Calibri"
NAVY = RGBColor(0x1F, 0x3A, 0x5C)
BLUE = RGBColor(0x2E, 0x86, 0xC1)
INK = RGBColor(0x26, 0x2B, 0x30)
MUTED = RGBColor(0x74, 0x7C, 0x84)
TILE_BG = RGBColor(0xEE, 0xF3, 0xF8)
TILE_LN = RGBColor(0xA8, 0xC2, 0xD6)
HILITE = RGBColor(0xDD, 0xEC, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
_num = [0]


# ---------------------------------------------------------------- low-level helpers
def _set(run, size, color, bold=False, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(26)
    return struct.unpack(">II", head[16:24])


def fit_image(slide, path, bl, bt, bw, bh):
    # Embeds the image scaled to fit the box, preserving aspect and centring.
    if not os.path.isfile(path):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(bl), Inches(bt), Inches(bw), Inches(bh))
        box.fill.solid(); box.fill.fore_color.rgb = TILE_BG
        box.line.color.rgb = TILE_LN; box.shadow.inherit = False
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _set(p.add_run(), 12, MUTED); p.runs[0].text = "missing: " + os.path.basename(path)
        return
    w, h = png_size(path)
    ar = w / h
    if ar > bw / bh:
        iw, ih = bw, bw / ar
    else:
        ih, iw = bh, bh * ar
    l = bl + (bw - iw) / 2
    t = bt + (bh - ih) / 2
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(iw), Inches(ih))


# ---------------------------------------------------------------- slide chrome
def chrome(slide, section, title):
    _num[0] += 1
    tf = textbox(slide, 8.4, 0.42, 4.33, 0.4)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _set(tf.paragraphs[0].add_run(), 12, BLUE, bold=True)
    tf.paragraphs[0].runs[0].text = section.upper()
    tf = textbox(slide, 0.6, 0.48, 8.0, 0.9, MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0].add_run(), 26, NAVY, bold=True)
    tf.paragraphs[0].runs[0].text = title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.4),
                                 Inches(2.2), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = textbox(slide, 12.35, 7.02, 0.8, 0.35)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _set(tf.paragraphs[0].add_run(), 9, MUTED)
    tf.paragraphs[0].runs[0].text = str(_num[0])


def content(section, title):
    s = prs.slides.add_slide(BLANK)
    chrome(s, section, title)
    return s


def bullets(slide, items, top=1.6, left=0.6, width=12.13, height=5.0):
    tf = textbox(slide, left, top, width, height)
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        prefix = "•  " if level == 0 else "–  "
        r = p.add_run(); r.text = prefix + text
        _set(r, 18 if level == 0 else 15, INK if level == 0 else MUTED)
    return tf


def caption(slide, text, top=6.9):
    tf = textbox(slide, 0.6, top, 12.13, 0.35)
    _set(tf.paragraphs[0].add_run(), 11, MUTED, italic=True)
    tf.paragraphs[0].runs[0].text = text


def col_label(slide, text, l, w, top=1.5):
    tf = textbox(slide, l, top, w, 0.4)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set(tf.paragraphs[0].add_run(), 15, BLUE, bold=True)
    tf.paragraphs[0].runs[0].text = text


def code_box(slide, lines, l, t, w, h, size=12):
    # Monospace card for real source snippets (code, prompt text, rubric excerpts).
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xF9)
    box.line.color.rgb = TILE_LN; box.line.width = Pt(0.75); box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_right = Pt(10)
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(1)
        r = p.add_run(); r.text = line if line else " "
        r.font.name = "Consolas"; r.font.size = Pt(size); r.font.color.rgb = INK


# ---------------------------------------------------------------- special slides
def divider(section_title):
    _num[0] += 1
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.95),
                             Inches(5.33), Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = textbox(s, 1.0, 2.9, 11.33, 1.1, MSO_ANCHOR.BOTTOM)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set(tf.paragraphs[0].add_run(), 40, NAVY, bold=True)
    tf.paragraphs[0].runs[0].text = section_title
    return s


def single_graph(section, title, img, note, cap=None):
    s = content(section, title)
    fit_image(s, img, 0.7, 1.55, 11.93, 5.15)
    if cap:
        caption(s, cap)
    notes(s, note)
    return s


def dual_graph(section, title, left_img, right_img, note, cap=None):
    s = content(section, title)
    labels = [("Qwen3:14b  ·  local", 0.5), ("gpt-oss-120b  ·  cloud", 6.83)]
    for text, x in labels:
        tf = textbox(s, x, 1.5, 6.0, 0.4)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set(tf.paragraphs[0].add_run(), 15, BLUE, bold=True)
        tf.paragraphs[0].runs[0].text = text
    fit_image(s, left_img, 0.5, 1.95, 6.0, 4.75)
    fit_image(s, right_img, 6.83, 1.95, 6.0, 4.75)
    if cap:
        caption(s, cap)
    notes(s, note)
    return s


# ================================================================ BUILD

# ---- Slide 1: title
s = prs.slides.add_slide(BLANK)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), SW, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
bar.line.fill.background(); bar.shadow.inherit = False
tf = textbox(s, 1.0, 1.35, 11.33, 1.2, MSO_ANCHOR.BOTTOM)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
_set(tf.paragraphs[0].add_run(), 34, NAVY, bold=True)
tf.paragraphs[0].runs[0].text = "Exploring Large Language Models for Grading Programming Assignments"
tf = textbox(s, 1.0, 2.7, 11.33, 0.6)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
_set(tf.paragraphs[0].add_run(), 18, MUTED, italic=True)
tf.paragraphs[0].runs[0].text = "A rubric-based, multi-approach comparison against human grades"
tf = textbox(s, 1.0, 3.9, 11.33, 1.4)
for i, line in enumerate([
    "MSc Computer Science (Intelligent Systems) — Dissertation",
    "[Your name]   ·   Supervisor: [name]   ·   [date]",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6)
    _set(p.add_run(), 15, INK if i == 0 else MUTED, bold=(i == 0))
    p.runs[0].text = line
notes(s, "Good [morning]. My dissertation asks a practical question: can large language "
      "models be trusted to grade student programming assignments — and if so, how should "
      "we set them up to do it? I'll take you through the motivation, the design, the "
      "implementation, the experiments on real and synthetic data, and the tool I built so "
      "a lecturer stays in control.")

# ================= MOTIVATION =================
divider("Motivation")

s = content("Motivation", "Grading code by hand doesn’t scale")
bullets(s, [
    (0, "Slow — hand-grading eats staff time as cohorts grow"),
    (0, "Subjective — two markers rarely agree on the same submission"),
    (0, "Doesn’t scale with class size"),
    (0, "Feedback shapes how students learn — quality matters, not just speed"),
    (0, "Pressure to automate — but it must be fair, not just fast"),
])
notes(s, "Grading code by hand is slow, subjective, and doesn’t scale as class sizes grow. "
      "And the feedback students get shapes how they learn. So there’s pressure to automate — "
      "the question is whether automation can do it fairly, not just fast.")

s = content("Motivation", "What today’s tools can — and can’t — do")
bullets(s, [
    (0, "Unit-test autograders (dominant approach)"),
    (1, "Messer et al. 2024 review of 121 tools: ~⅔ grade correctness only"),
    (1, "No partial credit, no open-ended reasoning; feedback is just pass/fail"),
    (0, "LLM-as-a-judge (emerging)"),
    (1, "Zheng 2023, Gu 2025: >80% agreement with humans on open-ended tasks"),
    (1, "Reads code, reasons, writes feedback → a plausible grader, not just a test runner"),
])
notes(s, "A 2024 review of 121 tools found two-thirds only check correctness, almost always "
      "with unit tests — which can’t grade partial credit, reasoning, or open-ended work. In "
      "parallel, LLM-as-a-judge shows >80% agreement with humans on open-ended tasks, making an "
      "LLM a plausible grader, not just a test runner.")

s = content("Motivation", "Can we trust it? The gap → six research questions")
bullets(s, [
    (0, "EU AI Act: automated grading is “high-risk” — human oversight mandated"),
    (0, "Students distrust opaque autograders"),
    (0, "Gap: prior work tests ONE variable at a time (one style / type / model / run)"),
    (0, "Aim: systematically explore reliable LLM grading → 6 RQs:"),
    (1, "RQ1 classifier vs free · RQ2 whole vs point · RQ3 robustness"),
    (1, "RQ4 bias · RQ5 feedback quality · RQ6 failure modes"),
])
notes(s, "Trust is the catch. The EU AI Act calls automated grading high-risk and mandates human "
      "oversight, and students distrust opaque autograders. The research gap: prior work tests one "
      "variable at a time. My aim is to explore reliable LLM grading systematically, broken into "
      "six research questions.")

# ================= DESIGN =================
divider("Design")

s = content("Design", "Two datasets: real + synthetic")
bullets(s, [
    (0, "Real — University of Hamburg “Algorithmen und Datenstrukturen” (AD2022)"),
    (1, "Anonymised student code, 3 winter terms, Java + Python (~1,526 solutions)"),
    (1, "Ecological validity — real, messy student code"),
    (0, "Synthetic — authored by me"),
    (1, "21 asymptotic-analysis questions (real corpus is implementation-only)"),
    (1, "Crafted correct / semi-correct / wrong solution variants"),
    (0, "Real proves it works in the wild; synthetic gives controlled coverage"),
])
notes(s, "Two complementary datasets. Real: the University of Hamburg Algorithms and Data "
      "Structures corpus — anonymised submissions over three winter terms, Java and Python. "
      "Synthetic, authored by me: 21 asymptotic-analysis questions (the real corpus is "
      "implementation-only), plus crafted correct/semi/wrong solution variants for controlled "
      "coverage. Real proves it works in the wild; synthetic controls exactly what I measure.")

s = content("Design", "Example question types")
col_label(s, "Implementation  ·  19_20-1-1-java", 0.5, 6.0)
col_label(s, "Asymptotic analysis  ·  asym-1-java", 6.83, 6.0)
tf = textbox(s, 0.5, 1.95, 6.0, 0.75)
_set(tf.paragraphs[0].add_run(), 12, INK, italic=True)
tf.paragraphs[0].runs[0].text = ('"Implement a function that calculates the LCM of two given '
                                 'values." (may not use import)')
tf.paragraphs[0].runs[0].font.size = Pt(12)
tf2 = textbox(s, 6.83, 1.95, 6.0, 0.75)
_set(tf2.paragraphs[0].add_run(), 12, INK, italic=True)
tf2.paragraphs[0].runs[0].text = ('"Calculate the asymptotic worst-case running time of f, as a '
                                  'function of N. Refer to line numbers."')
code_box(s, [
    "int gcd(int a, int b) {",
    "    if (a < b) {",
    "        int tmp = a; a = b; b = tmp;",
    "    }",
    "    int r = a % b;",
    "    if (r == 0) return b;",
    "    return gcd(r, b);",
    "}",
    "",
    "int lcm(int a, int b) {",
    "    return (a * b) / gcd(a, b);",
    "}",
], 0.5, 2.75, 6.0, 3.9, size=13)
code_box(s, [
    "static int f(int[] a) {",
    "    int n = a.length;",
    "    int s = 0;",
    "    for (int i = 0; i < n; i++) {",
    "        for (int j = 0; j <= i; j++) {",
    "            s += a[j] * a[i];",
    "        }",
    "    }",
    "    int t = g(a);",
    "    return s + t;",
    "}",
], 6.83, 2.75, 6.0, 3.9, size=13)
notes(s, "Two example questions: on the left, an implementation task from the real Hamburg corpus — "
      "compute the LCM using a GCD helper. On the right, one of my authored asymptotic-analysis "
      "questions — the student must derive the tight time-complexity bound and justify it by line "
      "number. Same rubric machinery grades both, even though one is 'does the code work' and the "
      "other is 'is the reasoning about the code correct'.")

s = content("Design", "Rubric design")
bullets(s, [
    (0, "Question-specific rubrics (beat generic ones — Pathak et al. 2025)"),
    (0, "~6 points per rubric; point marks sum to 100"),
    (0, "Exactly 3 buckets per point: Correct / Semi / Wrong, each with criteria"),
    (0, "Same rubric drives BOTH free-marking and classification — a fair comparison"),
])
notes(s, "Grading is only as good as the rubric. I use question-specific rubrics — Pathak et al. "
      "show these beat generic ones. Each has ~6 points; every point has exactly three buckets — "
      "Correct, Semi, Wrong — with criteria, and marks sum to 100. Same rubric drives both "
      "free-marking and classification, so the comparison is fair.")

s = content("Design", "Rubric example — one point, three buckets")
tf = textbox(s, 0.6, 1.55, 12.13, 0.6)
_set(tf.paragraphs[0].add_run(), 15, INK, bold=True)
tf.paragraphs[0].runs[0].text = "Point 1 · “gcd function present and structurally correct” — 15 marks (19_20-1-1)"
rt = s.shapes.add_table(2, 3, Inches(0.6), Inches(2.3), Inches(12.13), Inches(3.9)).table
rt.columns[0].width = Inches(4.04); rt.columns[1].width = Inches(4.04); rt.columns[2].width = Inches(4.05)
rt.rows[0].height = Inches(0.55)
WRONG_BG, WRONG_FG = RGBColor(0xFB, 0xE9, 0xE7), RGBColor(0xA6, 0x38, 0x2C)
SEMI_BG, SEMI_FG = RGBColor(0xFD, 0xF2, 0xDA), RGBColor(0xA9, 0x76, 0x0B)
CORRECT_BG, CORRECT_FG = RGBColor(0xE3, 0xF3, 0xE6), RGBColor(0x2E, 0x7D, 0x4F)
def _bucket_cell(c, text, size, bold, color, fill):
    c.margin_top = Pt(8); c.margin_bottom = Pt(8); c.margin_left = Pt(10); c.margin_right = Pt(10)
    c.vertical_anchor = MSO_ANCHOR.TOP
    c.fill.solid(); c.fill.fore_color.rgb = fill
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if bold else PP_ALIGN.LEFT
    _set(p.add_run(), size, color, bold=bold); p.runs[0].text = text
_bucket_cell(rt.cell(0, 0), "Wrong — 0 marks", 16, True, WRONG_FG, WRONG_BG)
_bucket_cell(rt.cell(0, 1), "Semi — 8 marks", 16, True, SEMI_FG, SEMI_BG)
_bucket_cell(rt.cell(0, 2), "Correct — 15 marks", 16, True, CORRECT_FG, CORRECT_BG)
_bucket_cell(rt.cell(1, 0),
             "No gcd function defined and no gcd-like logic present anywhere in the submission.",
             13, False, INK, WRONG_BG)
_bucket_cell(rt.cell(1, 1),
             "A gcd function is defined but has a broken signature or structure (e.g. missing "
             "return, wrong parameter count, no identifiable loop/recursion body).",
             13, False, INK, SEMI_BG)
_bucket_cell(rt.cell(1, 2),
             "A gcd function with the correct signature (two parameters, returns a value) and a "
             "recognisable recursive or iterative structure is present.",
             13, False, INK, CORRECT_BG)
caption(s, "Every rubric point uses this same 3-bucket structure — the same criteria drive both free-marking and classification.")
notes(s, "One concrete rubric point, in full: whether the student's gcd function is present and "
      "structurally correct, worth 15 marks. Wrong means no gcd logic at all; Semi means a gcd "
      "function exists but its signature or structure is broken; Correct means a properly-shaped "
      "gcd function is present — its actual correctness is checked by later points. Every point in "
      "every rubric follows this same three-bucket pattern.")

s = content("Design", "The four grading approaches (2×2)")
bullets(s, [(0, "Two axes: rubric granularity × scoring mode. Approaches 3 & 4 (classification) are the novel contribution.")],
        height=0.9)
gf = s.shapes.add_table(3, 3, Inches(2.4), Inches(2.9), Inches(8.5), Inches(2.6)).table
gf.columns[0].width = Inches(3.0); gf.columns[1].width = Inches(2.75); gf.columns[2].width = Inches(2.75)
def _cell(c, text, bold=False, color=INK, fill=WHITE, align=PP_ALIGN.LEFT):
    c.margin_top = Pt(3); c.margin_bottom = Pt(3); c.margin_left = Pt(6)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    c.fill.solid(); c.fill.fore_color.rgb = fill
    p = c.text_frame.paragraphs[0]; p.alignment = align
    _set(p.add_run(), 14, color, bold=bold); p.runs[0].text = text
_cell(gf.cell(0, 0), "")
_cell(gf.cell(0, 1), "Whole rubric", True, WHITE, NAVY, PP_ALIGN.CENTER)
_cell(gf.cell(0, 2), "Point-by-point", True, WHITE, NAVY, PP_ALIGN.CENTER)
_cell(gf.cell(1, 0), "Free marks", True, WHITE, NAVY)
_cell(gf.cell(2, 0), "Classify → bucket", True, WHITE, NAVY)
_cell(gf.cell(1, 1), "Approach 1", align=PP_ALIGN.CENTER)
_cell(gf.cell(1, 2), "Approach 2", align=PP_ALIGN.CENTER)
_cell(gf.cell(2, 1), "Approach 3", True, NAVY, HILITE, PP_ALIGN.CENTER)
_cell(gf.cell(2, 2), "Approach 4", True, NAVY, HILITE, PP_ALIGN.CENTER)
caption(s, "Every approach runs 3× per (question, student, model); median = final score.", top=5.7)
notes(s, "This 2×2 is the heart of the design. Two decisions: granularity (whole rubric vs one "
      "point at a time) and scoring mode (free marks vs classify into a fixed bucket). Approaches "
      "3 and 4 — grading as classification — is the novel part no prior work tests.")

s = content("Design", "Ground truth, repeated runs & metrics")
bullets(s, [
    (0, "Human ground truth (per-point buckets) = the yardstick for everything"),
    (0, "Every case run 3× → median (self-consistency, Wang et al. 2023)"),
    (0, "Metrics: MAE, Spearman ρ, Cohen’s κ, run std-dev, format-OK %, time"),
    (0, "Leniency = signed bias — catches hidden strictness/generosity correlation misses"),
])
notes(s, "I graded a reference set by hand, bucket by bucket — the yardstick for everything. "
      "Because LLMs are stochastic, every case runs three times and I take the median. Metrics: "
      "MAE and Spearman for accuracy and ranking, Cohen’s kappa for buckets, run std-dev for "
      "stability, format-OK, time — and leniency, the signed bias correlation is blind to.")

# ================= IMPLEMENTATION =================
divider("Implementation")

s = content("Implementation", "Pipeline architecture")
fit_image(s, ARCH, 0.5, 1.5, 12.33, 5.35)
notes(s, "A six-stage pipeline: preprocessing cleans the raw dataset into a datastore; validation "
      "checks it; the grade stage builds prompts and calls the models; processing parses responses "
      "into scores; metrics computes everything; and the last stage generates the plots. One "
      "orchestrator runs it end-to-end, and every stage is skippable so long runs resume.")

s = content("Implementation", "Prompting & multi-backend grading")
bullets(s, [
    (0, "Prompt = question + sample solution + student code + rubric (whole or one point)"),
    (0, "System prompt pins output to a strict, parseable format (POINT_ID / MARKS|BUCKET / FEEDBACK / TOTAL)"),
    (0, "Three backends, one code path:"),
    (1, "Ollama (local) · NVIDIA NIM (cloud) · generic OpenAI-compatible (Gemini, …)"),
    (0, "Swap any model without touching the grading logic"),
])
notes(s, "Each prompt gives the question, a reference solution, the student’s code, and the rubric "
      "— whole or single point — plus a system prompt that pins output to a strict parseable format. "
      "The same code talks to Ollama, NVIDIA NIM, and any OpenAI-compatible provider, so I can swap "
      "models without touching the grading logic.")

s = content("Implementation", "Prompt design — the four approaches")
tf = textbox(s, 0.6, 1.46, 12.13, 0.35)
_set(tf.paragraphs[0].add_run(), 12, MUTED, italic=True)
tf.paragraphs[0].runs[0].text = ("Same header every time (question + sample + student code). "
    "Rows = scoring mode · Columns = rubric granularity — only the instruction + output format change.")

def _prompt_card(label, lines, l, t):
    lab = textbox(s, l, t, 6.0, 0.32)
    _set(lab.paragraphs[0].add_run(), 13, BLUE, bold=True)
    lab.paragraphs[0].runs[0].text = label
    code_box(s, lines, l, t + 0.32, 6.0, 2.10, size=11)

A1 = [
    "Given: whole RUBRIC (all points + ranges)",
    "Assign a numeric mark per point.",
    "",
    "POINT_ID: <id>",
    "MARKS: <number>",
    "FEEDBACK: <1-2 sentences>",
    "  ... one block per point ...",
    "TOTAL: <sum of all MARKS>",
]
A2 = [
    "Given: ONE point  [one call per point]",
    "Assign a numeric mark for THIS point.",
    "",
    "MARKS: <number>",
    "FEEDBACK: <1-2 sentences>",
]
A3 = [
    "Given: whole RUBRIC (all points + buckets)",
    "Choose ONE bucket label per point.",
    "",
    "POINT_ID: <id>",
    "BUCKET: <exact bucket label>",
    "FEEDBACK: <1-2 sentences>",
    "  ... one block per point ...",
]
A4 = [
    "Given: ONE point  [one call per point]",
    "Choose ONE bucket label for THIS point.",
    "",
    "BUCKET: <exact bucket label>",
    "FEEDBACK: <1-2 sentences>",
]
_prompt_card("Approach 1  ·  Whole rubric + Free marks", A1, 0.5, 1.9)
_prompt_card("Approach 2  ·  Point-by-point + Free marks", A2, 6.83, 1.9)
_prompt_card("Approach 3  ·  Whole rubric + Classify", A3, 0.5, 4.5)
_prompt_card("Approach 4  ·  Point-by-point + Classify", A4, 6.83, 4.5)
notes(s, "All four prompts share the same header — question, sample solution, student code — and "
      "differ only in two axes. Left column (1 & 3) sends the whole rubric in one call; right column "
      "(2 & 4) sends one point per call. Top row (1 & 2) asks for a free numeric mark; bottom row "
      "(3 & 4) asks the model to classify into a bucket label. Only Approach 1 prints a TOTAL — and "
      "even then the pipeline ignores it and re-sums the per-point marks itself. This is the 2×2 "
      "design made visible in the prompts.")

s = content("Implementation", "Engineering realities")
bullets(s, [
    (0, "Robust response parser + a format-OK flag on every response"),
    (0, "Retries on timeout/empty; completed runs skipped → long runs restartable"),
    (0, "Fixed a real bug: a reasoning model exhausted its token budget “thinking”"),
    (1, "returned empty → disabled thinking-mode for that model"),
    (0, "These details are why the results are trustworthy, not noisy"),
])
notes(s, "Getting local models to reliably produce parseable output was real work: the parser flags "
      "format compliance, calls retry, completed cases are skipped so runs restart. One model kept "
      "burning its token budget ‘thinking’ and returned nothing — I disabled thinking-mode for it.")

# ================= EXPERIMENTS =================
divider("Experiments")

s = content("Experiments", "Two-phase experimental design")
bullets(s, [
    (0, "Phase 1 — Model selection (controlled, synthetic set)"),
    (1, "10 questions × correct/semi/wrong × 5 models × 4 approaches × 3 runs = 1,800 gradings"),
    (0, "Phase 2 — RQ study with the chosen models vs human grades"),
    (0, "First earn trust in a model, then use it to study how to grade"),
])
notes(s, "Two phases. First a controlled model-selection study on the synthetic set — 1,800 "
      "gradings — to earn the right to trust a model. Then the RQ study with the chosen models. "
      "First pick a trustworthy grader, then study how to grade.")

# ---- Model selection: text + 3 graphs
s = content("Experiments", "Phase 1 — Model selection")
bullets(s, [
    (0, "Choose the best local + cloud model from 5 candidates"),
    (0, "Judged on MAE, stability, format compliance, leniency, and speed"),
    (0, "Chosen: Qwen3:14b (local) + gpt-oss-120b (cloud)"),
    (0, "Next: full comparison table, then MAE / stability / format-failure graphs"),
])
notes(s, "Before the research questions I chose which models to trust, comparing five candidates on "
      "error, stability, format compliance, leniency and speed. The small models broke format; I "
      "picked Qwen3-14b locally and gpt-oss-120b in the cloud, prioritising accuracy and calibration. "
      "[→ read the deciding numbers from these graphs live.]")

# ---- Model comparison table (5 candidates)
s = content("Experiments · Model selection", "Model comparison — 5 candidates")
hdr = ["Model", "MAE ↓", "ρ ↑", "κ ↑", "Leniency →0", "Max dev ↓", "Format ↑", "Time(s) ↓"]
mrows = [
    ("devstral-small-2",  "6.65",  "0.942", "0.684", "−2.78", "2.28", "100.0%", "23.9",  False),
    ("Qwen3:14b",         "6.19",  "0.936", "0.735", "+0.78", "5.71", "98.6%",  "62.9",  True),
    ("gemma4:26b",        "8.89",  "0.857*","0.698*","−1.14", "8.65", "90.6%",  "106.3", False),
    ("nemotron-3-super-120b-a12b", "10.45", "0.830", "0.761", "−7.00", "7.58", "88.9%", "119.1", False),
    ("gpt-oss-120b",      "7.27",  "0.942", "0.728", "+0.78", "7.79", "94.2%",  "17.8",  True),
]
mt = s.shapes.add_table(6, 8, Inches(0.4), Inches(1.9), Inches(12.53), Inches(3.1)).table
mt.columns[0].width = Inches(3.15)
for c in range(1, 8):
    mt.columns[c].width = Inches(1.34)
for c, htxt in enumerate(hdr):
    cell = mt.cell(0, c)
    cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(6)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    _set(p.add_run(), 12, WHITE, bold=True); p.runs[0].text = htxt
for r, row in enumerate(mrows, start=1):
    chosen = row[8]
    fill = HILITE if chosen else (RGBColor(0xF5, 0xF5, 0xF5) if row[0] == "gemma4:26b" else WHITE)
    for c in range(8):
        cell = mt.cell(r, c)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.margin_left = Pt(6)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        color = NAVY if chosen else (MUTED if row[0] == "gemma4:26b" else INK)
        _set(p.add_run(), 12, color, bold=(chosen and c == 0)); p.runs[0].text = row[c]
bullets(s, [
    (0, "Highlighted rows = chosen models · greyed row = gemma4:26b, dropped"),
], top=5.15, height=0.5)
caption(s, "* gemma4’s ρ/κ computed on a pooled basis for this comparison, not identical methodology to the other four rows.", top=6.55)
notes(s, "Full picture across all five candidates. Devstral is the most stable and has zero format "
      "failures, but is stricter and less accurate. Nemotron has the best bucket agreement but is the "
      "strictest, slowest, and least format-compliant. Gemma4:26b — the largest model tried at 26B "
      "parameters — actually underperforms the 14B Qwen: worse accuracy, the worst run-to-run "
      "stability, and roughly 9% format failures on long prompts, which is why it was dropped despite "
      "its size. Qwen3-14b and gpt-oss-120b win on the combination of accuracy and calibration, which "
      "is why they're the two models carried into the RQ study.")

single_graph("Experiments · Model selection", "MAE by approach — 5 candidate models",
             os.path.join(MC, "mae_by_approach.png"),
             "Mean absolute error from human grades, per approach, for all five candidate models. "
             "[→ point out which models are most accurate.]")
single_graph("Experiments · Model selection", "Stability by model",
             os.path.join(MC, "stability_by_model.png"),
             "How much each model’s score moves across the three runs. [→ note the most and least "
             "stable models.]")
single_graph("Experiments · Model selection", "Format-failure count",
             os.path.join(MC, "format_failures.png"),
             "How often each model broke the required output format. [→ note which small/large models "
             "failed and why they were dropped.]")

# ---- RQ1
s = content("Experiments · RQ1", "RQ1 — Classifier vs free-marker")
bullets(s, [
    (0, "Question: classify each rubric point into a bucket, or let the LLM freely assign marks?"),
    (0, "Check: run-to-run consistency; agreement with human; does free-marking score lower?"),
    (0, "Novel — no prior work tests LLM-as-classifier for grading"),
])
notes(s, "RQ1: should the model classify into buckets or freely assign marks? Every graph here breaks "
      "down by approach, so the comparison is the classify approaches (A3/A4) against the free-mark ones "
      "(A1/A2), on run-to-run consistency and agreement with human. The classifier only picks among three "
      "buckets, so we’d expect less drift. [→ state the finding from the graphs.]")
dual_graph("Experiments · RQ1", "RQ1 · Score distribution vs human",
           os.path.join(QW, "score_distribution.png"), os.path.join(GO, "score_distribution.png"),
           "One score-distribution curve per approach plus the human curve, per model. "
           "[→ point out the classifier curves (A3/A4) vs the free-mark ones (A1/A2), and which sits closest to human.]",
           cap="One curve per approach + human — compare classifier (A3/A4) vs free-mark (A1/A2), and which is closest to human.")
dual_graph("Experiments · RQ1", "RQ1 · Run consistency (implementation)",
           os.path.join(QW, "runs_implementation.png"), os.path.join(GO, "runs_implementation.png"),
           "2×2 panels, one per approach, showing the 3 runs per solution variant on an implementation "
           "question. [→ does the classifier (A3/A4) hold steadier than free-marking (A1/A2)?]",
           cap="One panel per approach — does the classifier (A3/A4) hold steadier across the 3 runs than free-marking (A1/A2)?")
dual_graph("Experiments · RQ1", "RQ1 · Run consistency (asymptotic)",
           os.path.join(QW, "runs_asymptotic.png"), os.path.join(GO, "runs_asymptotic.png"),
           "Same per-approach run-stability view on an asymptotic-analysis question.",
           cap="Same per-approach view on an asymptotic question — classifier vs free-mark run stability.")
dual_graph("Experiments · RQ1", "RQ1 · Median deviation across runs",
           os.path.join(QW, "deviation_all_questions.png"), os.path.join(GO, "deviation_all_questions.png"),
           "Max deviation from the median run, by approach, across all questions — lower = steadier. "
           "[→ compare classify vs free.]",
           cap="Max run-to-run deviation by approach (all questions) — lower = steadier; classify (A3/A4) vs free (A1/A2).")

# ---- RQ2
s = content("Experiments · RQ2", "RQ2 — Whole rubric vs point-by-point")
bullets(s, [
    (0, "Question: pass the whole rubric at once, or one point per call?"),
    (0, "Check: which aligns better with human; time & cost (one call per rubric point)"),
    (0, "Pathak et al. 2025: point-wise tends to score lower"),
])
notes(s, "RQ2: whole rubric in one call, or one point at a time? I compare accuracy by approach and "
      "time. Point-by-point makes one call per rubric point, so it’s inherently costlier. [→ state "
      "whether it buys better agreement and whether it grades lower.]")
single_graph("Experiments · RQ2", "RQ2 · MAE by approach (both models)",
             os.path.join(QW, "mae_by_approach.png"),
             "MAE per approach for both models. [→ compare whole-rubric (A1/A3) vs point-by-point (A2/A4); "
             "does point-wise grade lower?]",
             cap="MAE per approach, both models — compare whole-rubric (A1/A3) vs point-by-point (A2/A4).")
dual_graph("Experiments · RQ2", "RQ2 · Spearman ρ by approach",
           os.path.join(QW, "spearman_by_approach.png"), os.path.join(GO, "spearman_by_approach.png"),
           "Rank correlation with human, grouped Whole rubric vs Per point. [→ which granularity agrees better.]",
           cap="Rank agreement with human, grouped Whole rubric (A1/A3) vs Per-point (A2/A4).")
single_graph("Experiments · RQ2", "RQ2 · Time by approach (both models)",
             os.path.join(QW, "time_by_approach.png"),
             "Grading time, grouped Whole rubric vs Per point. Point-by-point costs one call per rubric "
             "point. [→ quantify the accuracy-vs-compute trade-off.]",
             cap="Grading time, grouped Whole rubric vs Per-point — point-by-point = one call per rubric point.")

# ---- RQ3
s = content("Experiments · RQ3", "RQ3 — Robustness / stability")
bullets(s, [
    (0, "Question: how stable are marks across repeated runs? Does it depend on model & approach?"),
    (0, "Check: score variation across 3 runs; classification vs free; format-failure rate"),
    (0, "Schroeder 2025: a grader can be accurate yet unreliable"),
])
notes(s, "RQ3: how stable are marks across runs, and does it depend on model and approach? A grader "
      "can be accurate yet unreliable — giving different marks on identical input. [→ note which model "
      "and approach are most stable, and format-failure rates.]")
single_graph("Experiments · RQ3", "RQ3 · Stability by model & approach",
             os.path.join(QW, "stability_by_model.png"),
             "Boxplots of max deviation across the 3 runs, by model and approach. [→ tighter box = more "
             "robust; note the most stable model and whether classify beats free.]",
             cap="Max deviation across the 3 runs, by model and approach — tighter box = more robust.")
single_graph("Experiments · RQ3", "RQ3 · Format failures by model & approach",
             os.path.join(QW, "format_failures.png"),
             "Format-failure count by model and approach. [→ which model/approach broke format most.]",
             cap="Format-failure count by model and approach.")

# ---- RQ4
s = content("Experiments · RQ4", "RQ4 — Bias")
bullets(s, [
    (0, "Question: too strict or too generous? middle-bucket avoidance? length bias?"),
    (0, "Check: leniency; bucket distribution vs human; length ↔ marks correlation"),
    (0, "Ye 2024 “verbosity bias” · Jukiewicz 2025 liberal vs restrictive graders"),
])
notes(s, "RQ4: is the model biased — too strict or generous, avoiding the middle bucket, or rewarding "
      "longer code regardless of quality? Ye et al. call that verbosity bias. These plots let me measure "
      "where each model sits. [→ state the direction of any bias found.]")
dual_graph("Experiments · RQ4", "RQ4 · Leniency by approach",
           os.path.join(QW, "leniency_by_approach.png"), os.path.join(GO, "leniency_by_approach.png"),
           "Mean (LLM − human) by approach, per model. Positive = generous, negative = strict. "
           "[→ note the direction and whether classify vs free differ.]",
           cap="Mean (LLM − human) by approach — positive = generous, negative = strict.")
dual_graph("Experiments · RQ4", "RQ4 · Leniency per submission",
           os.path.join(QW, "leniency_scatter.png"), os.path.join(GO, "leniency_scatter.png"),
           "Per-submission (LLM − human) by approach, per model. [→ spread, and whether it systematically "
           "over- or under-marks.]",
           cap="Per-submission (LLM − human) by approach — spread and direction of over/under-marking.")
dual_graph("Experiments · RQ4", "RQ4 · Bucket distribution vs human",
           os.path.join(QW, "bucket_distribution.png"), os.path.join(GO, "bucket_distribution.png"),
           "Bucket pickup % (Wrong/Semi/Correct): human vs each approach. [→ does the model avoid the "
           "middle 'Semi' bucket relative to human?]",
           cap="Bucket pickup % — human vs each approach: does the model avoid the middle 'Semi' bucket?")
dual_graph("Experiments · RQ4", "RQ4 · Length bias",
           os.path.join(QW, "length_bias.png"), os.path.join(GO, "length_bias.png"),
           "Solution length vs median marks, one panel per approach, with Spearman ρ. [→ does longer code "
           "score higher regardless of quality?]",
           cap="Solution length vs marks, per approach, with ρ — does longer code score higher?")

# ---- RQ5
s = content("Experiments · RQ5", "RQ5 — Feedback quality")
bullets(s, [
    (0, "Question: is feedback grounded & consistent with the mark, or generic/hallucinated?"),
    (0, "Check: manual 0–4 score; automatic feedback-contradiction rate"),
    (0, "Jia 2024: 23–27% of LLM feedback hallucinated · Qian 2025: 16-dim framework"),
    (0, "Result is a table / manual analysis (no graph)"),
])
notes(s, "RQ5: is the feedback trustworthy — grounded in the student’s code and consistent with the "
      "mark — or generic and hallucinated? Jia et al. found ~a quarter of LLM feedback is hallucinated. "
      "I score grounding, factual correctness and score-feedback agreement, plus an automatic "
      "contradiction rate. [→ present the feedback table.]")

# ---- RQ6
s = content("Experiments · RQ6", "RQ6 — Failure modes")
bullets(s, [
    (0, "Question: what mistakes recur — are they systematic or random?"),
    (0, "Check: largest human–LLM gaps, each tagged with an error type"),
    (1, "hallucination · wrong rubric point · length-for-marks · code wrong but high marks"),
    (0, "Result is a tagged failure-case table (no graph)"),
])
notes(s, "RQ6: what kind of mistakes does the model make, and are they systematic? I take the largest "
      "human–LLM gaps and tag each with an error type. [→ present the tagged failure cases.] RQ5 and "
      "RQ6 are my most qualitative, still-maturing analyses.")

# ================= GUI =================
divider("GUI")

s = content("GUI", "Why a GUI: human-in-the-loop")
bullets(s, [
    (0, "LLM grades  →  Streamlit review tool  →  human finalises"),
    (0, "EU AI Act + SURE (Korthals 2026): human oversight; review only uncertain cases"),
    (1, "→ cuts human grading effort by 40–90%"),
    (0, "Deliverable is decision-support, NOT a replacement for the lecturer"),
])
notes(s, "Every serious paper — and the EU AI Act — says the model must not decide alone. Korthals et "
      "al. show reviewing only the uncertain cases cuts grading effort 40–90%. So the deliverable is "
      "decision-support that keeps the lecturer in control.")

s = content("GUI", "The review tool (live demo)")
bullets(s, [
    (0, "Overview & Outliers: clickable mark histogram, class stats, bulk “finalise all”"),
    (0, "Review & Adjust: student code + sample + rubric side by side"),
    (1, "all four approaches shown together; per-point bucket / marks / feedback editor"),
    (1, "pre-fill from any approach OR from human ground truth, then save final grade"),
    (0, "LLM does the first pass; the human keeps the final say"),
])
caption(s, "Live demo preferred; run: streamlit run src/explorer/overview.py")
notes(s, "The Overview page gives a clickable mark distribution with class stats and a bulk finalise. "
      "Review & Adjust opens one submission: question, student code, sample and rubric side by side; "
      "all four approaches together; a per-point editor; and pre-fill from any approach or from ground "
      "truth. The LLM does the first pass; the human keeps the final say.")

s = content("GUI", "Multi-provider & cohort view")
bullets(s, [
    (0, "Provider-agnostic — Ollama (local) / NVIDIA NIM / Gemini"),
    (0, "Rubric-point difficulty: which criteria students found hardest"),
    (0, "RQ plots on demand inside the tool"),
    (0, "Compare whole cohorts across assignments over time"),
])
notes(s, "The tool is provider-agnostic, surfaces which rubric points students found hardest, shows RQ "
      "plots on demand, and compares whole cohorts over time. A department could point it at any model "
      "and any course.")

# ================= CONCLUSION =================
divider("Conclusion")

s = content("Conclusion", "Contributions, limitations & future work")
bullets(s, [
    (0, "Contributions"),
    (1, "First joint comparison: classifier-vs-free × whole-vs-point, local + cloud, multi-run"),
    (1, "Implementation + asymptotic questions, real + synthetic data; a pipeline + human-in-the-loop tool"),
    (0, "What the results answer (headline from the final graphs): RQ1–RQ6"),
    (0, "Limitations: single-grader ground truth · authored asymptotic questions · feedback partly manual"),
    (0, "Future: scale real-student run · multi-grader truth · de-biasing"),
])
caption(s, "Thank you — questions?", top=6.6)
notes(s, "To close: the first study to compare classifier-vs-free-marking and whole-vs-point-by-point "
      "together, across local and cloud models, over repeated runs, on implementation and asymptotic "
      "questions, and on real and synthetic data — with a full pipeline and a human-in-the-loop tool. "
      "[→ deliver one headline per research question.] Limitations: single-grader ground truth, "
      "self-authored asymptotic questions, partly-manual feedback analysis. Thank you.")

prs.save(OUT)
print("Wrote", OUT, "with", len(prs.slides._sldIdLst), "slides")
